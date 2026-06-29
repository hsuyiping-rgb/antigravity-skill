"""
soil-image-deck 打包腳本

支援兩種模式：
- baked（預設）：圖裡已含文字，pptx 每頁一張 full-bleed 圖即可
- plate：圖為無文字底圖，依 YAML spec 疊加可編輯文字框
"""
import os
import sys
import argparse
import glob
from pathlib import Path
import yaml
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def hex_to_rgb(h: str) -> RGBColor:
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

def crop_to_ratio(src_path: str, target_w_in: float, target_h_in: float,
                   images_dir: Path = None) -> str:
    """依目標寬高比中央裁切，避免 python-pptx 強制拉伸變形。"""
    if not src_path:
        return src_path
    target_ratio = target_w_in / target_h_in
    img = Image.open(src_path)
    w, h = img.size
    current_ratio = w / h
    if abs(current_ratio - target_ratio) < 0.01:
        return src_path
    if current_ratio > target_ratio:
        new_w = int(h * target_ratio)
        left = (w - new_w) // 2
        img = img.crop((left, 0, left + new_w, h))
    else:
        new_h = int(w / target_ratio)
        top = (h - new_h) // 2
        img = img.crop((0, top, w, top + new_h))
    cache_dir = (images_dir or Path(src_path).parent) / "cropped"
    cache_dir.mkdir(parents=True, exist_ok=True)
    out = cache_dir / f"{Path(src_path).stem}__{target_w_in:.2f}x{target_h_in:.2f}.png"
    img.save(out)
    return str(out)

def run_baked(images_dir, output_file):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    
    # Sort files naturally
    img_files = sorted(glob.glob(os.path.join(images_dir, "*.png")) + glob.glob(os.path.join(images_dir, "*.jpg")))
    
    if not img_files:
        print(f"No images found in {images_dir}")
        sys.exit(1)
        
    for img_path in img_files:
        slide = prs.slides.add_slide(blank_layout)
        cropped_img = crop_to_ratio(img_path, 13.333, 7.5, Path(images_dir))
        slide.shapes.add_picture(cropped_img, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
        print(f"Added slide for {img_path}")
        
    prs.save(output_file)
    print(f"Successfully created {output_file} in baked mode.")

def run_plate(yaml_path, images_dir, output_file):
    with open(yaml_path, "r", encoding="utf-8") as f:
        spec = yaml.safe_load(f)
        
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    
    palette = spec.get("palette", {
        "bg": "#0D1B2A",
        "primary": "#00C6FF",
        "highlight": "#FFD700",
        "text": "#FFFFFF",
        "muted": "#A5B4CB"
    })
    
    default_font = spec.get("default_font", "Microsoft JhengHei")
    title_font = spec.get("title_font", "GenSekiGothic2 TW H")
    
    pages = spec.get("pages", [])
    for page in pages:
        slide = prs.slides.add_slide(blank_layout)
        
        # Add background image if available
        bg_name = page.get("background_image")
        if bg_name:
            bg_path = os.path.join(images_dir, bg_name)
            if os.path.exists(bg_path):
                cropped_img = crop_to_ratio(bg_path, 13.333, 7.5, Path(images_dir))
                slide.shapes.add_picture(cropped_img, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
            else:
                print(f"Warning: background image {bg_path} not found.")
                
        # Add textboxes
        textboxes = page.get("textboxes", [])
        for tb in textboxes:
            x = Inches(float(tb["x"]))
            y = Inches(float(tb["y"]))
            w = Inches(float(tb["w"]))
            h = Inches(float(tb["h"]))
            
            textbox = slide.shapes.add_textbox(x, y, w, h)
            tf = textbox.text_frame
            tf.word_wrap = True
            
            align_name = tb.get("align", "left").lower()
            if align_name == "center":
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            elif align_name == "right":
                tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
            else:
                tf.paragraphs[0].alignment = PP_ALIGN.LEFT
                
            p = tf.paragraphs[0]
            p.text = tb.get("text", "")
            
            font = p.font
            font.name = title_font if tb.get("is_title") else default_font
            font.size = Pt(int(tb.get("size", 20)))
            font.bold = bool(tb.get("bold", False))
            
            color_hex = tb.get("color", "text")
            if color_hex in palette:
                color_hex = palette[color_hex]
            font.color.rgb = hex_to_rgb(color_hex)
            
    prs.save(output_file)
    print(f"Successfully created {output_file} in plate mode.")

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--mode", default="baked", choices=["baked", "plate"])
    parser.add_argument("--yaml", default=None, help="YAML configuration file path for plate mode")
    parser.add_argument("--images-dir", required=True, help="Directory containing generated images")
    parser.add_argument("--output", default="output.pptx", help="Path to save the output PPTX file")
    args = parser.parse_args()
    
    if args.mode == "baked":
        run_baked(args.images_dir, args.output)
    elif args.mode == "plate":
        if not args.yaml:
            print("Error: --yaml is required for plate mode.")
            sys.exit(1)
        run_plate(args.yaml, args.images_dir, args.output)

if __name__ == "__main__":
    main()
